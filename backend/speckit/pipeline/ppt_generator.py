from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os
from typing import Dict, Any
from datetime import datetime
import re

class VAPowerPointGenerator:
    def __init__(self):
        # Modern VA Color Palette - Professional Medical Theme
        self.colors = {
            'primary': RGBColor(0, 51, 102),      # Deep Navy Blue
            'secondary': RGBColor(0, 102, 204),   # VA Blue
            'accent': RGBColor(14, 78, 128),      # Teal Accent
            'highlight': RGBColor(236, 112, 99),  # Coral for emphasis
            'light': RGBColor(248, 249, 250),     # Off-white
            'gray': RGBColor(108, 117, 125),      # Medium Gray
            'dark_gray': RGBColor(52, 58, 64),    # Dark Gray
            'white': RGBColor(255, 255, 255),     # Pure White
            'success': RGBColor(40, 167, 69),     # Success Green
            'warning': RGBColor(255, 193, 7),     # Warning Yellow
            'info': RGBColor(23, 162, 184),       # Info Cyan
        }
        
        self.icons_dir = os.path.join(os.path.dirname(__file__), '..', '..', 'assets', 'icons')
        self.logos_dir = os.path.join(os.path.dirname(__file__), '..', '..', 'assets', 'logos')
        
        # Specialty keywords for auto-detection
        self.specialty_keywords = {
            'cardiology': ['cardio', 'heart', 'blood pressure', 'hypertension', 'myocardial', 'cardiac', 'cardiovascular', 'coronary'],
            'neurology': ['neuro', 'brain', 'stroke', 'cognitive', 'dementia', 'alzheimer', 'parkinson', 'seizure'],
            'oncology': ['cancer', 'tumor', 'oncology', 'metastasis', 'chemotherapy', 'radiation', 'malignant'],
            'infectious_disease': ['infection', 'viral', 'bacterial', 'covid', 'sepsis', 'antibiotic', 'pathogen'],
            'surgery': ['surgery', 'operative', 'surgical', 'procedure', 'laparoscopic'],
            'pharmacy': ['drug', 'medication', 'pharmac', 'therapeutic', 'dosage'],
            'pulmonology': ['lung', 'respiratory', 'pneumonia', 'asthma', 'copd', 'pulmonary'],
            'psychiatry': ['depress', 'anxiety', 'psychiat', 'mental', 'psychological'],
            'orthopedics': ['bone', 'fracture', 'orthop', 'arthritis', 'joint', 'musculoskeletal'],
            'dermatology': ['skin', 'derma', 'rash', 'cutaneous'],
            'gastroenterology': ['GI', 'gastro', 'liver', 'intestinal', 'digestive', 'bowel'],
        }
    
    def generate_presentation(self, summaries: Dict[str, Any], medical_icon: str, job_id: str) -> Dict[str, Any]:
        """Create modern, icon-rich VA presentation following professional design principles."""
        try:
            # Auto-detect specialty
            if not medical_icon or medical_icon == 'general':
                medical_icon = self._detect_specialty(summaries) or 'general_medicine'
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # 1. Title slide with specialty icon
            self._add_professional_title_slide(prs, summaries, medical_icon)
            
            # 2. Study Design & Methods slide
            self._add_study_design_slide(prs, summaries, medical_icon)
            
            # 3. Key Findings with visual emphasis
            self._add_visual_findings_slide(prs, summaries, medical_icon)
            
            # 4. Results breakdown (if numeric data available)
            self._add_results_slide(prs, summaries, medical_icon)
            
            # 5. Clinical Implications & Conclusion
            self._add_implications_slide(prs, summaries, medical_icon)
            
            # Save
            output_dir = 'output'
            os.makedirs(output_dir, exist_ok=True)
            filename = f"va_abstract_{job_id}.pptx"
            file_path = os.path.join(output_dir, filename)
            prs.save(file_path)
            
            return {
                'success': True,
                'file_path': file_path,
                'file_size': os.path.getsize(file_path),
                'filename': filename,
                'specialty': medical_icon
            }
        except Exception as e:
            return {'success': False, 'message': str(e), 'error_type': 'ppt_generation_error'}
    
    def _add_professional_title_slide(self, prs, summaries, medical_icon):
        """Professional title slide with large specialty icon."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Clean white background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['white']
        bg.line.fill.background()
        
        # Top accent bar with VA blue
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.colors['primary']
        accent.line.fill.background()
        
        # VA Logo (smaller, top-left)
        logo_path = os.path.join(self.logos_dir, 'va_logo.png')
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.4), width=Inches(1))
        
        # Large Specialty Icon (center-top)
        icon_path = os.path.join(self.icons_dir, f'{medical_icon}.png')
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, 
                                    prs.slide_width/2 - Inches(1), 
                                    Inches(1.2), 
                                    width=Inches(2))
        
        # Title with professional typography
        title = summaries.get('title', 'Clinical Study Abstract')
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), prs.slide_width - Inches(2), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.colors['primary']
        
        # Subtitle with study details
        population = summaries.get('population', '')
        setting = summaries.get('setting', '')
        if population or setting:
            sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), prs.slide_width - Inches(3), Inches(1))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            subtitle_parts = [p for p in [population, setting] if p]
            sp.text = " • ".join(subtitle_parts)
            sp.alignment = PP_ALIGN.CENTER
            sp.font.name = 'Calibri'
            sp.font.size = Pt(18)
            sp.font.color.rgb = self.colors['gray']
        
        # Bottom date/attribution
        footer = slide.shapes.add_textbox(0, prs.slide_height - Inches(0.6), prs.slide_width, Inches(0.4))
        fp = footer.text_frame.paragraphs[0]
        fp.text = f"Visual Abstract • {datetime.now().strftime('%B %Y')}"
        fp.alignment = PP_ALIGN.CENTER
        fp.font.name = 'Calibri'
        fp.font.size = Pt(12)
        fp.font.color.rgb = self.colors['gray']
    
    def _add_study_design_slide(self, prs, summaries, medical_icon):
        """Study design slide with icons and structured layout."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['light']
        bg.line.fill.background()
        
        # Header with icon
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.colors['primary']
        header_bar.line.fill.background()
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(8), Inches(0.5))
        tp = title.text_frame.paragraphs[0]
        tp.text = 'STUDY DESIGN & METHODS'
        tp.font.name = 'Calibri'
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = self.colors['white']
        
        # Small specialty icon in header
        icon_path = os.path.join(self.icons_dir, f'{medical_icon}.png')
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, prs.slide_width - Inches(1.3), Inches(0.15), width=Inches(0.7))
        
        # Create information cards with icons
        cards = [
            ('👥 POPULATION', summaries.get('population', 'Not specified'), self.colors['secondary']),
            ('💊 INTERVENTION', summaries.get('intervention', 'Not specified'), self.colors['info']),
            ('🏥 SETTING', summaries.get('setting', 'Not specified'), self.colors['success']),
            ('🎯 PRIMARY OUTCOME', summaries.get('primary_outcome', 'Not specified'), self.colors['highlight']),
        ]
        
        card_width = (prs.slide_width - Inches(2)) / 2
        card_height = Inches(2.2)
        
        for i, (label, text, color) in enumerate(cards):
            row = i // 2
            col = i % 2
            left = Inches(0.6) + col * (card_width + Inches(0.8))
            top = Inches(1.5) + row * (card_height + Inches(0.4))
            
            # Card with colored left border
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors['white']
            card.line.color.rgb = color
            card.line.width = Pt(3)
            
            # Colored accent bar on left
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.15), card_height)
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = color
            accent_bar.line.fill.background()
            
            # Card content
            tf = card.text_frame
            tf.margin_left = Inches(0.3)
            tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.25)
            tf.word_wrap = True
            
            # Label with emoji
            lp = tf.paragraphs[0]
            lp.text = label
            lp.font.name = 'Calibri'
            lp.font.size = Pt(14)
            lp.font.bold = True
            lp.font.color.rgb = color
            lp.space_after = Pt(8)
            
            # Text content
            tp = tf.add_paragraph()
            tp.text = text[:200]  # Truncate if too long
            tp.font.name = 'Calibri'
            tp.font.size = Pt(13)
            tp.font.color.rgb = self.colors['dark_gray']
            tp.space_after = Pt(0)
    
    def _add_visual_findings_slide(self, prs, summaries, medical_icon):
        """Key findings with visual emphasis and statistical highlights."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['white']
        bg.line.fill.background()
        
        # Colored header bar
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.colors['success']
        header_bar.line.fill.background()
        
        # Title
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(8), Inches(0.5))
        tp = title.text_frame.paragraphs[0]
        tp.text = '✓ KEY FINDINGS'
        tp.font.name = 'Calibri'
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = self.colors['white']
        
        # Extract statistics from findings
        findings_text = summaries.get('findings', 'No findings available.')
        stats = self._extract_statistics(findings_text)
        
        # Statistics highlight boxes (if any)
        if stats:
            stat_y = Inches(1.3)
            stat_width = (prs.slide_width - Inches(2)) / min(len(stats), 3)
            for i, stat in enumerate(stats[:3]):  # Max 3 stats
                left = Inches(0.7) + i * stat_width
                self._add_stat_box(slide, left, stat_y, stat_width - Inches(0.2), stat)
        
        # Main findings content box
        findings_top = Inches(2.8) if stats else Inches(1.3)
        findings_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Inches(0.7), findings_top,
                                             prs.slide_width - Inches(1.4), 
                                             prs.slide_height - findings_top - Inches(0.7))
        findings_box.fill.solid()
        findings_box.fill.fore_color.rgb = RGBColor(245, 247, 250)
        findings_box.line.color.rgb = self.colors['secondary']
        findings_box.line.width = Pt(2)
        
        # Findings bullet points
        tf = findings_box.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.4)
        tf.word_wrap = True
        
        # Split findings into bullets
        bullets = [s.strip() for s in findings_text.split('.') if s.strip() and len(s.strip()) > 10][:6]
        
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f'▸ {bullet}.'
            p.font.name = 'Calibri'
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark_gray']
            p.space_after = Pt(14)
            p.level = 0
    
    def _add_results_slide(self, prs, summaries, medical_icon):
        """Results slide with data visualization elements."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light gradient background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['light']
        bg.line.fill.background()
        
        # Header
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.colors['info']
        header_bar.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(8), Inches(0.5))
        tp = title.text_frame.paragraphs[0]
        tp.text = '📊 RESULTS'
        tp.font.name = 'Calibri'
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = self.colors['white']
        
        # Primary Outcome
        outcome_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(0.7), Inches(1.4),
                                            prs.slide_width - Inches(1.4), Inches(1.5))
        outcome_box.fill.solid()
        outcome_box.fill.fore_color.rgb = self.colors['white']
        outcome_box.line.color.rgb = self.colors['primary']
        outcome_box.line.width = Pt(3)
        
        otf = outcome_box.text_frame
        otf.margin_left = Inches(0.4)
        otf.margin_top = Inches(0.25)
        otf.word_wrap = True
        
        op = otf.paragraphs[0]
        op.text = 'PRIMARY OUTCOME'
        op.font.name = 'Calibri'
        op.font.size = Pt(14)
        op.font.bold = True
        op.font.color.rgb = self.colors['secondary']
        
        outcome_text = summaries.get('primary_outcome', 'Not specified')
        op2 = otf.add_paragraph()
        op2.text = outcome_text
        op2.font.name = 'Calibri'
        op2.font.size = Pt(18)
        op2.font.color.rgb = self.colors['dark_gray']
        
        # Secondary findings
        secondary_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Inches(0.7), Inches(3.2),
                                              prs.slide_width - Inches(1.4), 
                                              prs.slide_height - Inches(4))
        secondary_box.fill.solid()
        secondary_box.fill.fore_color.rgb = RGBColor(250, 252, 254)
        secondary_box.line.color.rgb = self.colors['secondary']
        secondary_box.line.width = Pt(1.5)
        
        stf = secondary_box.text_frame
        stf.margin_left = Inches(0.4)
        stf.margin_top = Inches(0.3)
        stf.word_wrap = True
        
        sp = stf.paragraphs[0]
        sp.text = 'SECONDARY FINDINGS'
        sp.font.name = 'Calibri'
        sp.font.size = Pt(14)
        sp.font.bold = True
        sp.font.color.rgb = self.colors['secondary']
        sp.space_after = Pt(10)
        
        # Extract additional findings
        findings = summaries.get('findings', '')
        secondary_points = [s.strip() for s in findings.split('.') if s.strip() and len(s.strip()) > 15][1:4]
        
        for point in secondary_points:
            sp2 = stf.add_paragraph()
            sp2.text = f'• {point}.'
            sp2.font.name = 'Calibri'
            sp2.font.size = Pt(14)
            sp2.font.color.rgb = self.colors['dark_gray']
            sp2.space_after = Pt(8)
    
    def _add_implications_slide(self, prs, summaries, medical_icon):
        """Clinical implications and conclusion slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['white']
        bg.line.fill.background()
        
        # Header with gradient effect
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.colors['primary']
        header_bar.line.fill.background()
        
        title = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(8), Inches(0.5))
        tp = title.text_frame.paragraphs[0]
        tp.text = '💡 CLINICAL IMPLICATIONS'
        tp.font.name = 'Calibri'
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = self.colors['white']
        
        # Large specialty icon (watermark style)
        icon_path = os.path.join(self.icons_dir, f'{medical_icon}.png')
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, 
                                    prs.slide_width - Inches(3), 
                                    prs.slide_height - Inches(3.5), 
                                    width=Inches(2.5))
        
        # Main conclusion box
        conclusion_text = summaries.get('va_summary') or summaries.get('conclusion', 'No conclusion provided.')
        
        conclusion_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               Inches(0.7), Inches(1.5),
                                               prs.slide_width - Inches(4.5), 
                                               prs.slide_height - Inches(2.5))
        conclusion_box.fill.solid()
        conclusion_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
        conclusion_box.line.color.rgb = self.colors['primary']
        conclusion_box.line.width = Pt(4)
        
        tf = conclusion_box.text_frame
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.5)
        tf.margin_top = Inches(0.5)
        tf.word_wrap = True
        
        # Main message
        p = tf.paragraphs[0]
        p.text = conclusion_text
        p.font.name = 'Calibri'
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors['dark_gray']
        p.space_after = Pt(20)
        p.line_spacing = 1.3
        
        # Footer with attribution
        footer = slide.shapes.add_textbox(0, prs.slide_height - Inches(0.6), prs.slide_width, Inches(0.4))
        fp = footer.text_frame.paragraphs[0]
        fp.text = f"Visual Abstract Generated by JAMA VA Abstractor • {datetime.now().strftime('%B %Y')}"
        fp.alignment = PP_ALIGN.CENTER
        fp.font.name = 'Calibri'
        fp.font.size = Pt(11)
        fp.font.color.rgb = self.colors['gray']
    
    def _add_stat_box(self, slide, left, top, width, stat_text):
        """Add a highlighted statistics box."""
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = self.colors['secondary']
        box.line.fill.background()
        
        tf = box.text_frame
        tf.vertical_anchor = 1  # Middle
        p = tf.paragraphs[0]
        p.text = stat_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
    
    def _extract_statistics(self, text):
        """Extract statistical values from text."""
        stats = []
        # Look for patterns like "X%", "p<0.05", "95% CI", numbers with units
        patterns = [
            r'\d+\.?\d*%',  # Percentages
            r'p\s*[<>=]\s*0\.\d+',  # P-values
            r'\d+\.?\d*\s*mmHg',  # Blood pressure
            r'HR[:\s]+\d+\.?\d*',  # Hazard ratio
            r'OR[:\s]+\d+\.?\d*',  # Odds ratio
            r'RR[:\s]+\d+\.?\d*',  # Relative risk
        ]
        
        for pattern in patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            stats.extend(matches[:3])  # Limit to avoid cluttering
            if len(stats) >= 3:
                break
        
        return stats[:3]
    
    def _detect_specialty(self, summaries):
        """Detect specialty from text."""
        text = ' '.join([
            str(summaries.get(k, '')).lower() 
            for k in ('title', 'findings', 'intervention', 'population')
        ])
        
        for specialty, keywords in self.specialty_keywords.items():
            if any(kw in text for kw in keywords):
                return specialty
        
        return 'general_medicine'

# Test
if __name__ == '__main__':
    test_summaries = {
        'title': 'Digital Health Interventions for Cardiovascular Risk Reduction',
        'population': '500 veterans aged 50-75 with hypertension',
        'intervention': 'Mobile app with daily BP monitoring',
        'setting': 'VA community health centers nationwide',
        'primary_outcome': 'Systolic BP reduction at 12 weeks',
        'findings': 'Mean 8.5 mmHg reduction (95% CI: 6.2-10.8, p<0.001). Significant improvement in BP control. High user engagement (85% daily use). Reduced cardiovascular risk markers.',
        'va_summary': 'Mobile health intervention significantly reduces blood pressure in veteran population with excellent user engagement and clinical benefit.'
    }
    
    gen = VAPowerPointGenerator()
    result = gen.generate_presentation(test_summaries, '', 'modern_demo')
    print(f"{'✓' if result['success'] else '✗'} {result.get('file_path', result.get('message'))}")
