from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
from typing import Dict, Any
from datetime import datetime
from io import BytesIO

# Optional charting support
try:
    import matplotlib.pyplot as plt
    HAS_MATPLOTLIB = True
except Exception:
    HAS_MATPLOTLIB = False

# Optional Pillow support for icon generation
try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except Exception:
    HAS_PIL = False


class VAPowerPointGenerator:
    def __init__(self):
        # VA Color Scheme
        self.va_colors = {
            'navy': RGBColor(0, 32, 96),
            'accent': RGBColor(14, 78, 128),
            'teal': RGBColor(0, 120, 135),
            'gray': RGBColor(105, 105, 105),
            'muted': RGBColor(240, 242, 245),
            'white': RGBColor(255, 255, 255),
            'green': RGBColor(34, 139, 34),
            'red': RGBColor(220, 20, 60),
        }

        self.icon_symbols = {
            'cardiology': '♥',
            'neurology': '🧠',
            'oncology': '※',
            'infectious_disease': '⚡',
            'surgery': '✚',
            'pharmacy': '⚕',
            'general_medicine': '⚕'
        }

        # icons directory (optional PNG icons per specialty)
        self.icons_dir = os.path.join(os.path.dirname(__file__), '..', '..', 'assets', 'icons')

        # Keyword mapping for specialty detection
        self.specialty_keywords = {
            'cardiology': ['cardio', 'heart', 'blood pressure', 'hypertension', 'myocardial', 'cardiac', 'arrhythmia'],
            'neurology': ['neuro', 'brain', 'stroke', 'neurolog', 'seizure', 'cognitive', 'dementia'],
            'oncology': ['cancer', 'tumor', 'oncology', 'neoplasm', 'metastasis'],
            'infectious_disease': ['infection', 'viral', 'bacterial', 'covid', 'sepsis', 'infectious'],
            'surgery': ['surgery', 'operative', 'procedure', 'lapar', 'surgical'],
            'pharmacy': ['drug', 'medication', 'pharmac', 'dose', 'pharmacy'],
            'pulmonology': ['lung', 'respiratory', 'pneumonia', 'asthma', 'copd'],
            'psychiatry': ['depress', 'anxiety', 'psychiat', 'mental', 'suicide'],
            'orthopedics': ['bone', 'fracture', 'orthop', 'hip', 'knee', 'arthritis'],
            'dermatology': ['skin', 'derma', 'rash', 'psoriasis', 'eczema'],
            'gastroenterology': ['GI', 'gastro', 'liver', 'intestinal', 'colitis', 'abdomen']
        }

    def generate_presentation(self, summaries: Dict[str, Any], medical_icon: str, job_id: str) -> Dict[str, Any]:
        """Create a multi-slide, polished VA-style presentation."""
        try:
            # auto-detect specialty icon if not provided or generic
            if not medical_icon or medical_icon == 'general' or medical_icon not in self.icon_symbols:
                detected = self._detect_specialty(summaries)
                medical_icon = detected or 'general_medicine'

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Title slide
            self._add_title_slide(prs, summaries, medical_icon)

            # Overview slide (study details)
            self._add_overview_slide(prs, summaries)

            # Findings + KPIs + optional chart
            self._add_findings_slide(prs, summaries)

            # Methods slide (if available)
            if summaries.get('methods') or summaries.get('study_design'):
                self._add_methods_slide(prs, summaries)

            # Results slide (detailed results or tables)
            if summaries.get('results') or summaries.get('detailed_results'):
                self._add_results_slide(prs, summaries)

            # Conclusion slide
            self._add_conclusion_slide(prs, summaries)

            # Save
            output_dir = 'output'
            os.makedirs(output_dir, exist_ok=True)
            filename = f"va_abstract_{job_id}.pptx"
            file_path = os.path.join(output_dir, filename)
            prs.save(file_path)

            return {
                'success': True,
                'file_path': file_path,
                'file_size': os.path.getsize(file_path),
                'filename': filename
            }

        except Exception as e:
            return {'success': False, 'message': str(e), 'error_type': 'ppt_generation_error'}

    def _add_title_slide(self, prs: Presentation, summaries: Dict[str, Any], medical_icon: str):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.va_colors['navy']
        bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), prs.slide_width - Inches(1.2), Inches(1.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = summaries.get('title', 'Clinical Study Abstract')
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.size = Pt(32)
        r.font.bold = True
        r.font.name = 'Calibri'
        r.font.color.rgb = self.va_colors['navy']

        # Subtitle with population/intervention
        subtitle_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), prs.slide_width - Inches(2.4), Inches(0.8))
        stf = subtitle_box.text_frame
        sp = stf.paragraphs[0]
        population = summaries.get('population', '')
        intervention = summaries.get('intervention', '')
        sp.text = f"{population} — {intervention}" if population or intervention else summaries.get('setting', '')
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.runs[0]
        sr.font.size = Pt(14)
        sr.font.color.rgb = self.va_colors['gray']
        sr.font.name = 'Calibri'

        # Icon: try PNG asset, then Pillow-generated circle, else emoji
        icon_path = os.path.join(self.icons_dir, f"{medical_icon}.png")
        if os.path.exists(icon_path):
            try:
                slide.shapes.add_picture(icon_path, Inches(1.0), Inches(0.9), width=Inches(1.4), height=Inches(1.4))
            except Exception:
                # fallback to emoji text
                self._add_icon_emoji(slide, medical_icon)
        elif HAS_PIL:
            # generate simple circular icon with initial
            try:
                buf = self._generate_icon_image(medical_icon)
                slide.shapes.add_picture(buf, Inches(1.0), Inches(0.9), width=Inches(1.4), height=Inches(1.4))
            except Exception:
                self._add_icon_emoji(slide, medical_icon)
        else:
            self._add_icon_emoji(slide, medical_icon)

    def _add_overview_slide(self, prs: Presentation, summaries: Dict[str, Any]):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        left = Inches(0.6)
        top = Inches(0.5)
        col_width = (prs.slide_width - Inches(1.2)) / 2

        # Left: Study details as cards
        cards = [
            ('Population', summaries.get('population', 'N/A')),
            ('Intervention', summaries.get('intervention', 'N/A')),
            ('Setting', summaries.get('setting', 'N/A')),
            ('Primary Outcome', summaries.get('primary_outcome', 'N/A')),
        ]

        card_h = Inches(1.0)
        for i, (label, text) in enumerate(cards):
            top_i = top + (card_h + Inches(0.15)) * i
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_i, col_width - Inches(0.15), card_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.va_colors['muted']
            shape.line.color.rgb = self.va_colors['navy']
            shape.line.width = Pt(1)

            tf = shape.text_frame
            tp = tf.paragraphs[0]
            tp.text = label
            tp.runs[0].font.size = Pt(12)
            tp.runs[0].font.bold = True
            tp.runs[0].font.color.rgb = self.va_colors['navy']

            p2 = tf.add_paragraph()
            p2.text = text
            p2.runs[0].font.size = Pt(11)
            p2.runs[0].font.color.rgb = self.va_colors['gray']

        # Right: short list of study highlights
        right_left = left + col_width + Inches(0.3)
        title_box = slide.shapes.add_textbox(right_left, top, col_width, Inches(0.5))
        ttf = title_box.text_frame
        t = ttf.paragraphs[0]
        t.text = 'Study Highlights'
        t.runs[0].font.size = Pt(18)
        t.runs[0].font.bold = True
        t.runs[0].font.color.rgb = self.va_colors['navy']

        highlights = summaries.get('highlights') or [x for x in summaries.get('findings', '').split('.') if x][:5]
        for i, hl in enumerate(highlights[:5]):
            htop = top + Inches(0.5) + (Pt(18).pt/72 + 0.12) * i + Inches(0.05)
            tb = slide.shapes.add_textbox(right_left, top + Inches(0.5) + Inches(0.4)*i, col_width, Inches(0.4))
            tfb = tb.text_frame
            pp = tfb.paragraphs[0]
            pp.text = f'• {hl.strip()}'
            pp.runs[0].font.size = Pt(12)
            pp.runs[0].font.color.rgb = self.va_colors['gray']

    def _add_findings_slide(self, prs: Presentation, summaries: Dict[str, Any]):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        left = Inches(0.6)
        top = Inches(0.6)
        right = prs.slide_width - Inches(0.6)

        # Big findings box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, right - left, Inches(3.6))
        box.fill.solid()
        box.fill.fore_color.rgb = self.va_colors['white']
        box.line.color.rgb = self.va_colors['accent']
        box.line.width = Pt(2)

        tf = box.text_frame
        h = tf.paragraphs[0]
        h.text = 'Key Findings'
        h.runs[0].font.size = Pt(18)
        h.runs[0].font.bold = True
        h.runs[0].font.color.rgb = self.va_colors['navy']

        findings = summaries.get('findings', '')
        bullets = [s.strip() for s in findings.split('.') if s.strip()][:6]
        for b in bullets:
            p = tf.add_paragraph()
            p.text = f'• {b}'
            p.level = 1
            p.runs[0].font.size = Pt(13)
            p.runs[0].font.color.rgb = self.va_colors['gray']

        # KPI boxes under findings
        kpi_top = top + Inches(3.8)
        kpi_w = (right - left - Inches(0.8)) / 3
        kpis = []
        # Example: allow 'primary_effect' or 'quality_score' keys
        if summaries.get('primary_effect'):
            kpis.append(('Effect', summaries.get('primary_effect')))
        if summaries.get('quality_score') is not None:
            kpis.append(('Quality', f"{summaries.get('quality_score'):.0%}"))
        if summaries.get('n'):
            kpis.append(('N', str(summaries.get('n'))))

        for i, (label, value) in enumerate(kpis[:3]):
            left_i = left + Inches(0.2) + (kpi_w + Inches(0.2)) * i
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_i, kpi_top, kpi_w, Inches(0.9))
            shp.fill.solid()
            shp.fill.fore_color.rgb = self.va_colors['accent']
            shp.line.fill.background()
            tfk = shp.text_frame
            p1 = tfk.paragraphs[0]
            p1.text = label
            p1.runs[0].font.size = Pt(11)
            p1.runs[0].font.color.rgb = self.va_colors['white']
            p1.runs[0].font.bold = True

            p2 = tfk.add_paragraph()
            p2.text = value
            p2.runs[0].font.size = Pt(16)
            p2.runs[0].font.color.rgb = self.va_colors['white']
            p2.runs[0].font.bold = True

        # Optional chart
        if HAS_MATPLOTLIB and summaries.get('chart_data'):
            try:
                img = self._create_chart_image(summaries['chart_data'])
                pic = slide.shapes.add_picture(img, right - Inches(4.2), kpi_top, width=Inches(3.8))
            except Exception:
                pass

    def _add_icon_emoji(self, slide, medical_icon: str):
        icon_symbol = self.icon_symbols.get(medical_icon, self.icon_symbols['general_medicine'])
        icon_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.9), Inches(1.4), Inches(1.4))
        itf = icon_box.text_frame
        ip = itf.paragraphs[0]
        ip.text = icon_symbol
        ip.alignment = PP_ALIGN.CENTER
        ir = ip.runs[0]
        ir.font.size = Pt(36)
        ir.font.color.rgb = self.va_colors['white']

    def _generate_icon_image(self, key: str) -> BytesIO:
        """Create a simple circular PNG icon with an emoji/initial and return BytesIO."""
        size = (256, 256)
        bg = (14, 78, 128)
        fg = (255, 255, 255)
        img = Image.new('RGBA', size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)
        draw.ellipse((0, 0, size[0], size[1]), fill=bg)

        symbol = self.icon_symbols.get(key, self.icon_symbols['general_medicine'])
        try:
            # try emoji font first, fallback to default
            font = ImageFont.truetype('seguiemj.ttf', 160)
        except Exception:
            try:
                font = ImageFont.truetype('arial.ttf', 140)
            except Exception:
                font = ImageFont.load_default()

        w, h = draw.textsize(symbol, font=font)
        draw.text(((size[0]-w)/2, (size[1]-h)/2), symbol, font=font, fill=fg)

        buf = BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return buf

    def _detect_specialty(self, summaries: Dict[str, Any]) -> str:
        """Simple keyword-based specialty detection from title/findings/intervention."""
        text = ' '.join([
            str(summaries.get(k, '')).lower() for k in ('title', 'findings', 'intervention', 'population', 'setting')
        ])

        for specialty, keywords in self.specialty_keywords.items():
            for kw in keywords:
                if kw.lower() in text:
                    return specialty

        return 'general_medicine'

        # Footer
        self._add_footer(slide, prs.slide_width, prs.slide_height)

    def _add_methods_slide(self, prs: Presentation, summaries: Dict[str, Any]):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        left = Inches(0.6)
        top = Inches(0.6)
        width = prs.slide_width - Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        t = title_box.text_frame.paragraphs[0]
        t.text = 'Methods'
        t.runs[0].font.size = Pt(20)
        t.runs[0].font.bold = True
        t.runs[0].font.color.rgb = self.va_colors['navy']

        content = summaries.get('methods') or summaries.get('study_design') or 'Methods not provided.'
        cb = slide.shapes.add_textbox(left, top + Inches(0.9), width, Inches(4.6))
        cf = cb.text_frame
        for para in str(content).split('\n')[:8]:
            p = cf.add_paragraph() if cf.paragraphs else cf.paragraphs[0]
            p.text = para.strip()
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = self.va_colors['gray']

        self._add_footer(slide, prs.slide_width, prs.slide_height)

    def _add_results_slide(self, prs: Presentation, summaries: Dict[str, Any]):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        left = Inches(0.6)
        top = Inches(0.6)
        width = prs.slide_width - Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        t = title_box.text_frame.paragraphs[0]
        t.text = 'Results'
        t.runs[0].font.size = Pt(20)
        t.runs[0].font.bold = True
        t.runs[0].font.color.rgb = self.va_colors['navy']

        results = summaries.get('results') or summaries.get('detailed_results') or summaries.get('findings', '')
        rb = slide.shapes.add_textbox(left, top + Inches(0.9), width, Inches(4.6))
        rf = rb.text_frame
        bullets = [s.strip() for s in str(results).split('.') if s.strip()][:10]
        for b in bullets:
            p = rf.add_paragraph()
            p.text = f'• {b}'
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = self.va_colors['gray']

        self._add_footer(slide, prs.slide_width, prs.slide_height)

    def _add_conclusion_slide(self, prs: Presentation, summaries: Dict[str, Any]):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        left = Inches(0.6)
        top = Inches(0.6)
        width = prs.slide_width - Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.6))
        t = title_box.text_frame.paragraphs[0]
        t.text = 'Conclusion & Implications'
        t.runs[0].font.size = Pt(20)
        t.runs[0].font.bold = True
        t.runs[0].font.color.rgb = self.va_colors['navy']

        conclusion = summaries.get('conclusion') or summaries.get('va_summary') or 'No conclusion provided.'
        cb = slide.shapes.add_textbox(left, top + Inches(0.9), width, Inches(3.6))
        cf = cb.text_frame
        p = cf.paragraphs[0]
        p.text = conclusion
        p.runs[0].font.size = Pt(13)
        p.runs[0].font.color.rgb = self.va_colors['gray']

        # Actionable recommendations (if provided)
        recs = summaries.get('recommendations') or summaries.get('implications')
        if recs:
            rb = slide.shapes.add_textbox(left, top + Inches(2.6), width, Inches(2.0))
            rf = rb.text_frame
            rf.paragraphs[0].text = 'Recommendations'
            rf.paragraphs[0].runs[0].font.bold = True
            for r in str(recs).split('\n')[:5]:
                p = rf.add_paragraph()
                p.text = f'• {r.strip()}'
                p.runs[0].font.size = Pt(12)
                p.runs[0].font.color.rgb = self.va_colors['gray']

        self._add_footer(slide, prs.slide_width, prs.slide_height)

    def _create_chart_image(self, chart_data: Dict[str, Any]) -> BytesIO:
        """Create a simple bar chart image and return a BytesIO object."""
        fig, ax = plt.subplots(figsize=(4, 2.5), dpi=100)
        labels = [d.get('label') for d in chart_data]
        values = [d.get('value') for d in chart_data]
        ax.bar(labels, values, color=['#0E4E80'])
        ax.set_ylabel(chart_data[0].get('ylabel', ''))
        ax.set_title(chart_data[0].get('title', ''))
        ax.grid(axis='y', linestyle='--', alpha=0.4)
        buf = BytesIO()
        plt.tight_layout()
        fig.savefig(buf, format='png', dpi=150, transparent=True)
        plt.close(fig)
        buf.seek(0)
        return buf

    def _add_footer(self, slide, slide_width, slide_height):
        footer_top = slide_height - Inches(0.5)
        footer_box = slide.shapes.add_textbox(Inches(0.5), footer_top, slide_width - Inches(1), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_p = footer_frame.paragraphs[0]
        current_date = datetime.now().strftime('%B %d, %Y')
        footer_p.text = f"Generated by JAMA VA Abstractor • {current_date}"
        footer_p.alignment = PP_ALIGN.CENTER
        footer_p.runs[0].font.size = Pt(10)
        footer_p.runs[0].font.color.rgb = self.va_colors['gray']


# Quick local test
if __name__ == '__main__':
    test_summaries = {
        'title': 'Digital Health Interventions for Cardiovascular Risk',
        'population': '500 veterans aged 50-75 with hypertension',
        'intervention': 'Mobile app with daily BP monitoring and coaching',
        'setting': 'VA community health centers nationwide',
        'primary_outcome': 'Systolic blood pressure reduction at 12 weeks',
        'findings': 'Mean 8.5 mmHg reduction (95% CI: 6.2-10.8, p=0.003). Clinically significant improvement in cardiovascular risk profile.',
        'va_summary': 'Mobile health intervention significantly reduces BP in veteran population with high user engagement and clinical benefit.',
        'primary_effect': '−8.5 mmHg',
        'quality_score': 0.87,
        'n': 500,
    }

    # optional chart data (will only show if matplotlib available)
    if HAS_MATPLOTLIB:
        test_summaries['chart_data'] = [
            {'label': 'Baseline', 'value': 142, 'title': 'Systolic BP', 'ylabel': 'mmHg'},
            {'label': '12w', 'value': 133}
        ]

    generator = VAPowerPointGenerator()
    res = generator.generate_presentation(test_summaries, 'cardiology', 'demo')
    if res['success']:
        print('Generated:', res['file_path'])
    else:
        print('Error:', res.get('message'))